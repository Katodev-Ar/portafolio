"""
PrintBridge Native Converter
Convierte DOCX, XLSX, PPTX, TXT, RTF → PDF usando solo Python.
Sin LibreOffice, sin dependencias externas.

Librerías usadas:
  - python-docx   → leer DOCX
  - openpyxl      → leer XLSX
  - python-pptx   → leer PPTX
  - reportlab     → generar PDF
  - Pillow        → imágenes embebidas
"""

import io
import os
import re
import base64
import logging
import textwrap
from pathlib import Path
from functools import lru_cache

log = logging.getLogger("PrintBridge.Converter")

# ── ReportLab imports ──────────────────────────────────────────────────────────
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, HRFlowable, KeepTogether, Image as RLImage,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4, letter, legal
from reportlab.lib import colors
from reportlab.lib.units import mm, inch
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.pdfbase import pdfmetrics
from PIL import Image as PILImage

# ── Page sizes ─────────────────────────────────────────────────────────────────
PAGE_SIZES = {
    "a4": A4, "letter": letter, "legal": legal,
}

# ── Conversion limits (B-01: constants previously missing → NameError) ─────────
MAX_ROWS = 500   # max rows rendered in XLSX → PDF conversion
MAX_COLS = 30    # max columns rendered in XLSX → PDF conversion

def _get_size(page_size):
    """Accept either a string key or a direct (w,h) tuple."""
    if isinstance(page_size, tuple):
        return page_size
    return PAGE_SIZES.get(str(page_size).lower(), A4)

# ─────────────────────────────────────────────────────────────────────────────
#  DOCX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def docx_to_pdf(docx_path: str, pdf_path: str, page_size="a4"):
    """Convert .docx to PDF using python-docx + reportlab."""
    from docx import Document
    from docx.shared import RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    doc = Document(docx_path)
    size = _get_size(page_size)
    story = []
    styles = _make_styles(size)

    def rgb_to_hex(rgb):
        if rgb is None:
            return "#000000"
        try:
            return "#{:02x}{:02x}{:02x}".format(rgb.red, rgb.green, rgb.blue)
        except Exception:
            return "#000000"

    def run_to_markup(run):
        """Convert a docx Run to ReportLab XML markup string."""
        text = run.text or ""
        if not text:
            return ""
        # Escape XML special chars (Fix MEDIO: incluir comillas dobles)
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")
        # Replace tabs with spaces
        text = text.replace("\t", "&nbsp;&nbsp;&nbsp;&nbsp;")

        tags_open, tags_close = [], []

        if run.bold:
            tags_open.append("<b>"); tags_close.insert(0, "</b>")
        if run.italic:
            tags_open.append("<i>"); tags_close.insert(0, "</i>")
        if run.underline:
            tags_open.append("<u>"); tags_close.insert(0, "</u>")

        # Font size
        fs = None
        try:
            if run.font.size:
                fs = int(run.font.size.pt)
        except Exception:
            pass

        # Font color
        fc = None
        try:
            if run.font.color and run.font.color.type is not None:
                fc = rgb_to_hex(run.font.color.rgb)
        except Exception:
            pass

        # Build font tag
        font_attrs = []
        if fs:
            font_attrs.append(f'size="{fs}"')
        if fc and fc != "#000000":
            font_attrs.append(f'color="{fc}"')
        if font_attrs:
            tags_open.insert(0, f'<font {" ".join(font_attrs)}>')
            tags_close.append("</font>")

        return "".join(tags_open) + text + "".join(tags_close)

    def para_style(para):
        """Pick the closest ReportLab style for a docx paragraph."""
        try:
            name = (para.style.name or "").lower() if para.style is not None else ""
        except Exception:
            name = ""
        align_map = {
            WD_ALIGN_PARAGRAPH.CENTER: TA_CENTER,
            WD_ALIGN_PARAGRAPH.RIGHT: TA_RIGHT,
            WD_ALIGN_PARAGRAPH.JUSTIFY: TA_JUSTIFY,
        }
        align = align_map.get(para.alignment, TA_LEFT)

        if "heading 1" in name:
            s = styles["h1"]
        elif "heading 2" in name:
            s = styles["h2"]
        elif "heading 3" in name:
            s = styles["h3"]
        elif "heading" in name:
            s = styles["h4"]
        elif "title" in name:
            s = styles["title"]
        elif "list" in name or "bullet" in name:
            s = styles["bullet"]
        else:
            s = styles["normal"]

        if align != TA_LEFT:
            s = ParagraphStyle(
                s.name + "_aligned", parent=s, alignment=align
            )
        return s

    # ── Image extraction helper ───────────────────────────────────────────────
    def extract_inline_images(para_element):
        """Extract all inline images from a paragraph element, return RLImage list."""
        WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
        A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
        R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        imgs = []
        # Find all blip references (inline images)
        for blip in para_element.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blip"):
            r_embed = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if not r_embed:
                continue
            try:
                part = doc.part.related_parts.get(r_embed)
                if part is None:
                    continue
                img_bytes = part.blob
                pil_img = PILImage.open(io.BytesIO(img_bytes))
                if pil_img.mode in ("RGBA", "P", "LA"):
                    pil_img = pil_img.convert("RGB")
                # Get desired size from extent element
                extent = para_element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}extent")
                max_w = 400  # default max width in points
                if extent is not None:
                    cx = int(extent.get("cx", 0))  # EMU
                    cy = int(extent.get("cy", 0))
                    if cx > 0:
                        max_w = min(cx / 914400 * 72, 450)  # EMU→points, cap at 450pt
                # Scale keeping aspect ratio
                iw, ih = pil_img.size
                scale = max_w / iw if iw > 0 else 1
                rh = ih * scale
                buf = io.BytesIO()
                pil_img.save(buf, "JPEG", quality=88)
                buf.seek(0)
                rl_img = RLImage(buf, width=max_w, height=rh)
                imgs.append(rl_img)
            except Exception as e:
                log.debug(f"DOCX image extraction error: {e}")
        return imgs

    # ── Process paragraphs and tables in document order ──────────────────────
    def process_body(body):
        for child in body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                from docx.text.paragraph import Paragraph as DocxPara
                para = DocxPara(child, doc)

                # Check for inline images first
                inline_imgs = extract_inline_images(child)
                if inline_imgs:
                    for img in inline_imgs:
                        story.append(img)
                        story.append(Spacer(1, 4))
                    # Also add any text in the same paragraph
                    markup = "".join(run_to_markup(r) for r in para.runs)
                    if markup.strip():
                        story.append(Paragraph(markup, para_style(para)))
                    continue

                markup = "".join(run_to_markup(r) for r in para.runs)

                if not markup.strip():
                    story.append(Spacer(1, 4))
                    continue

                # Numbering / bullets
                numPr = child.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}numPr")
                if numPr is not None:
                    ilvl_el = numPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ilvl")
                    level = int(ilvl_el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", 0)) if ilvl_el is not None else 0
                    indent = 12 + level * 12
                    bullet_style = ParagraphStyle(
                        "bullet_level", parent=styles["bullet"],
                        leftIndent=indent, bulletIndent=indent - 10,
                    )
                    story.append(Paragraph(f"• {markup}", bullet_style))
                else:
                    s = para_style(para)
                    story.append(Paragraph(markup, s))

            elif tag == "tbl":
                from docx.table import Table as DocxTable
                tbl = DocxTable(child, doc)
                _add_docx_table(tbl, story, styles)

            elif tag == "sectPr":
                pass  # Section properties — ignore

    process_body(doc.element.body)

    # ── Build PDF ─────────────────────────────────────────────────────────────
    doc_pdf = SimpleDocTemplate(
        pdf_path, pagesize=size,
        leftMargin=20*mm, rightMargin=20*mm,
        topMargin=20*mm, bottomMargin=20*mm,
    )
    doc_pdf.build(story)
    log.info(f"DOCX→PDF: {docx_path} → {pdf_path}")


def _add_docx_table(tbl, story, styles):
    """Convert a python-docx Table to a ReportLab Table."""
    data = []
    for row in tbl.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text or ""
            text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            row_data.append(Paragraph(text, styles["table_cell"]))
        data.append(row_data)

    if not data:
        return

    col_count = max(len(r) for r in data)
    # Equalize row lengths
    for r in data:
        while len(r) < col_count:
            r.append(Paragraph("", styles["table_cell"]))

    t = Table(data, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#3a3a3a")),
        ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",   (0, 0), (-1, 0), 9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f5f5")]),
        ("FONTSIZE",   (0, 1), (-1, -1), 8),
        ("GRID",       (0, 0), (-1, -1), 0.4, colors.HexColor("#cccccc")),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING",  (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("VALIGN",     (0, 0), (-1, -1), "TOP"),
    ]))
    story.append(t)
    story.append(Spacer(1, 8))


# ─────────────────────────────────────────────────────────────────────────────
#  XLSX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def xlsx_to_pdf(xlsx_path: str, pdf_path: str, page_size="a4", scale="fit"):
    """Convert .xlsx to PDF using openpyxl + reportlab."""
    import openpyxl
    from openpyxl.utils import get_column_letter

    # Fix Mejora #3: read_only=True evita cargar estilos/validaciones en memoria
    # Un Excel de 50 MB puede consumir 500 MB+ sin este flag
    wb = openpyxl.load_workbook(xlsx_path, data_only=True, read_only=True)
    try:
        size = _get_size(page_size)
        # page_size already resolved (landscape if needed)
        page = size

        story = []
        styles = _make_styles(page)

        for sheet_idx, ws in enumerate(wb.worksheets):
            if sheet_idx > 0:
                story.append(PageBreak())

            # Sheet title
            story.append(Paragraph(
                ws.title,
                ParagraphStyle("sheet_title", fontSize=13, fontName="Helvetica-Bold",
                               textColor=colors.HexColor("#1a1a2e"), spaceAfter=8)
            ))

            # Gather used range — en modo read_only, iter_rows es la forma correcta
            # max_column/max_row pueden ser None en read_only si el writer no los guardó
            max_col = min(ws.max_column or 1, 30)   # cap 30 cols
            max_row = min(ws.max_row or 1, 500)     # cap 500 rows

            # R1 Fix: usar iter_rows con límites explícitos — ws.cell() no existe en ReadOnlyWorksheet
            all_rows = list(ws.iter_rows(
                min_row=1, max_row=max_row,
                min_col=1, max_col=max_col,
                values_only=False
            ))
            if not all_rows:
                story.append(Paragraph("(hoja vacía)", styles["normal"]))
                continue

            # Build table data
            data = []
            for ri, row_cells in enumerate(all_rows, start=1):
                row_data = []
                for ci, cell in enumerate(row_cells, start=1):
                    val = cell.value
                    if val is None:
                        val = ""
                    text = str(val)
                    # Escape XML
                    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

                    # Style: first row = header
                    if ri == 1:
                        p = Paragraph(text, ParagraphStyle(
                            "xhdr", fontSize=7, fontName="Helvetica-Bold",
                            textColor=colors.white, leading=9
                        ))
                    else:
                        p = Paragraph(text, ParagraphStyle(
                            "xcell", fontSize=7, fontName="Helvetica",
                            textColor=colors.black, leading=9
                        ))
                    row_data.append(p)
                data.append(row_data)

            if not data:
                continue

            # Calculate column widths (proportional to content)
            avail_w = page[0] - 30 * mm
            col_w = avail_w / max_col
            col_widths = [col_w] * max_col

            t = Table(data, colWidths=col_widths, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND",    (0, 0), (-1, 0), colors.HexColor("#2d6a4f")),
                ("TEXTCOLOR",     (0, 0), (-1, 0), colors.white),
                ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f0f7f4")]),
                ("GRID",          (0, 0), (-1, -1), 0.3, colors.HexColor("#cccccc")),
                ("TOPPADDING",    (0, 0), (-1, -1), 2),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
                ("LEFTPADDING",   (0, 0), (-1, -1), 3),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 3),
                ("VALIGN",        (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(t)

            if max_row >= MAX_ROWS:
                story.append(Spacer(1, 6))
                story.append(Paragraph(
                    f"⚠ Vista limitada a {MAX_ROWS} filas.",
                    ParagraphStyle("warn", fontSize=7, textColor=colors.HexColor("#cc6600"))
                ))

        doc_pdf = SimpleDocTemplate(
            pdf_path, pagesize=page,
            leftMargin=15*mm, rightMargin=15*mm,
            topMargin=15*mm, bottomMargin=15*mm,
        )
        doc_pdf.build(story)
        log.info(f"XLSX→PDF: {xlsx_path} → {pdf_path}")
    finally:
        # Fix MEDIO: cerrar siempre el workbook
        wb.close()


# ─────────────────────────────────────────────────────────────────────────────
#  PPTX → PDF
# ─────────────────────────────────────────────────────────────────────────────

def pptx_to_pdf(pptx_path: str, pdf_path: str, page_size="a4"):
    """
    Convert .pptx to PDF by rendering each slide as an image via Pillow,
    then assembling into a PDF with reportlab.
    Each slide is drawn as a canvas with background, shapes, and text.
    """
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptxRGB

    prs = Presentation(pptx_path)

    # Slide dimensions
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height
    slide_w_px = int(slide_w_emu / 914400 * 96)   # EMU → 96dpi pixels
    slide_h_px = int(slide_h_emu / 914400 * 96)

    # Scale up for quality
    SCALE = 2.5
    W = int(slide_w_px * SCALE)
    H = int(slide_h_px * SCALE)

    def emu_to_px(emu):
        return int(emu / 914400 * 96 * SCALE)

    def rgb_from_pptx(color_obj):
        """Extract (r,g,b) tuple from pptx color."""
        try:
            if color_obj and color_obj.type is not None:
                c = color_obj.rgb
                return (c.red, c.green, c.blue)
        except Exception:
            pass
        return None

    def get_bg_color(slide):
        """Get slide background color."""
        try:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            c = rgb_from_pptx(fill.fore_color)
            if c:
                return c
        except Exception:
            pass
        return (255, 255, 255)

    # Build PDF with reportlab using images
    size = _get_size(page_size)
    pdf_canvas = rl_canvas.Canvas(pdf_path, pagesize=size)  # NEW-04
    page_w, page_h = size

    for slide_num, slide in enumerate(prs.slides):
        img = PILImage.new("RGB", (W, H), get_bg_color(slide))
        draw = PILImage.new("RGBA", (W, H), (0, 0, 0, 0))

        from PIL import ImageDraw, ImageFont

        imd = ImageDraw.Draw(img)

        # Load fonts — P-03: lru_cache para evitar recargar fuentes en cada run
        @lru_cache(maxsize=64)
        def get_font(size_pt, bold=False, italic=False):
            size_px = max(8, int(size_pt * SCALE * 96 / 72))
            font_names = []
            if bold and italic:
                font_names = ["arialbi.ttf", "DejaVuSans-BoldOblique.ttf"]
            elif bold:
                font_names = ["arialbd.ttf", "DejaVuSans-Bold.ttf"]
            elif italic:
                font_names = ["ariali.ttf", "DejaVuSans-Oblique.ttf"]
            else:
                font_names = ["arial.ttf", "DejaVuSans.ttf"]
            for fn in font_names:
                try:
                    return ImageFont.truetype(fn, size_px)
                except Exception:
                    pass
            try:
                return ImageFont.load_default()
            except Exception:
                return None

        # Draw shapes
        for shape in slide.shapes:
            try:
                x = emu_to_px(shape.left or 0)
                y = emu_to_px(shape.top or 0)
                w = emu_to_px(shape.width or 0)
                h = emu_to_px(shape.height or 0)

                # Background fill for shape
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        fc = rgb_from_pptx(fill.fore_color)
                        if fc:
                            imd.rectangle([x, y, x+w, y+h], fill=fc)
                except Exception:
                    pass

                # Image shapes
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_blob = shape.image.blob
                        pil_img = PILImage.open(io.BytesIO(img_blob)).convert("RGBA")
                        pil_img = pil_img.resize((max(1, w), max(1, h)), PILImage.LANCZOS)
                        rgb_img = PILImage.new("RGB", pil_img.size, (255, 255, 255))
                        rgb_img.paste(pil_img, mask=pil_img.split()[3] if pil_img.mode == "RGBA" else None)
                        img.paste(rgb_img, (x, y))
                    except Exception as e:
                        log.debug(f"PPTX image shape error: {e}")

                # Text frames
                if shape.has_text_frame:
                    tf = shape.text_frame
                    ty = y + int(3 * SCALE)
                    for para in tf.paragraphs:
                        line_parts = []
                        max_fs = 14
                        for run in para.runs:
                            fs = 14
                            try:
                                if run.font.size:
                                    fs = int(run.font.size.pt)
                            except Exception:
                                pass
                            max_fs = max(max_fs, fs)
                            bold = False
                            italic = False
                            try:
                                bold = bool(run.font.bold)
                                italic = bool(run.font.italic)
                            except Exception:
                                pass
                            fc = (0, 0, 0)
                            try:
                                c2 = rgb_from_pptx(run.font.color)
                                if c2:
                                    fc = c2
                            except Exception:
                                pass
                            line_parts.append((run.text or "", fs, bold, italic, fc))

                        if not line_parts and para.text:
                            line_parts = [(para.text, 14, False, False, (0, 0, 0))]

                        # Draw each run on this line
                        tx = x + int(4 * SCALE)
                        for text, fs, bold, italic, fc in line_parts:
                            if not text:
                                continue
                            font = get_font(fs, bold, italic)
                            if font is None:
                                continue
                            # Wrap text within shape width
                            lines = _wrap_text_pil(imd, text, font, w - int(8 * SCALE))
                            for line in lines:
                                if ty + int(fs * SCALE * 1.3) > y + h:
                                    break
                                try:
                                    imd.text((tx, ty), line, fill=fc, font=font)
                                except Exception:
                                    pass
                                ty += int(fs * SCALE * 1.4)
                        ty += int(max_fs * SCALE * 0.3)  # paragraph spacing

            except Exception as e:
                log.debug(f"PPTX shape error: {e}")

        # Add slide to PDF
        buf = io.BytesIO()
        img.save(buf, "JPEG", quality=90)
        buf.seek(0)

        # Scale image to fit page
        img_ratio = W / H
        page_ratio = page_w / page_h
        margin = 10 * mm
        if img_ratio > page_ratio:
            draw_w = page_w - 2 * margin
            draw_h = draw_w / img_ratio
        else:
            draw_h = page_h - 2 * margin
            draw_w = draw_h * img_ratio

        x_off = (page_w - draw_w) / 2
        y_off = (page_h - draw_h) / 2

        from reportlab.lib.utils import ImageReader
        pdf_canvas.drawImage(ImageReader(buf), x_off, y_off, width=draw_w, height=draw_h)

        # Slide number
        pdf_canvas.setFont("Helvetica", 7)
        pdf_canvas.setFillColor(colors.grey)
        pdf_canvas.drawCentredString(page_w / 2, 6 * mm, f"{slide_num + 1} / {len(prs.slides)}")

        pdf_canvas.showPage()

    pdf_canvas.save()
    log.info(f"PPTX→PDF: {pptx_path} → {pdf_path}")


def _wrap_text_pil(draw, text, font, max_width):
    """Wrap text to fit within max_width pixels."""
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        try:
            bbox = draw.textbbox((0, 0), test, font=font)
            tw = bbox[2] - bbox[0]
        except Exception:
            tw = len(test) * 7  # fallback estimate
        if tw <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [text]


# ─────────────────────────────────────────────────────────────────────────────
#  TXT / RTF → PDF
# ─────────────────────────────────────────────────────────────────────────────

def txt_to_pdf(txt_path: str, pdf_path: str, page_size="a4"):
    """Convert plain text to PDF."""
    size = _get_size(page_size)
    with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    story = []
    styles = _make_styles(size)

    for line in text.splitlines():
        line_esc = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if line_esc.strip():
            story.append(Paragraph(line_esc, styles["mono"]))
        else:
            story.append(Spacer(1, 4))

    doc_pdf = SimpleDocTemplate(
        pdf_path, pagesize=size,
        leftMargin=20*mm, rightMargin=20*mm,
        topMargin=20*mm, bottomMargin=20*mm,
    )
    doc_pdf.build(story)
    log.info(f"TXT→PDF: {txt_path} → {pdf_path}")


def rtf_to_pdf(rtf_path: str, pdf_path: str, page_size="a4"):
    """
    Convert RTF to PDF by stripping RTF control codes and treating as text.
    Basic but functional for most RTF files.
    """
    with open(rtf_path, "r", encoding="utf-8", errors="replace") as f:
        rtf = f.read()

    # Strip RTF control words and groups
    text = _strip_rtf(rtf)
    tmp_txt = rtf_path + ".tmp.txt"
    with open(tmp_txt, "w", encoding="utf-8") as f:
        f.write(text)
    txt_to_pdf(tmp_txt, pdf_path, page_size)
    os.remove(tmp_txt)
    log.info(f"RTF→PDF: {rtf_path} → {pdf_path}")


def _strip_rtf(rtf: str) -> str:
    """
    RTF stripper que maneja anidación recursiva de grupos {}.
    Fix Bug #6: el regex anterior solo eliminaba grupos sin braces anidados,
    produciendo PDFs vacíos con la mayoría de archivos RTF reales.
    """
    # Eliminar grupos recursivamente hasta que no queden cambios
    prev = None
    text = rtf
    _MAX_RTF_ITER = 1000  # P-04: evitar O(N²) en RTF profundamente anidado
    for _ in range(_MAX_RTF_ITER):
        if prev == text:
            break
        prev = text
        text = re.sub(r'\{[^{}]*\}', '', text)
    else:
        log.warning("_strip_rtf: hit max iterations — RTF may be deeply nested")
    # Remove control words and sequences
    text = re.sub(r'\\[a-z]+\-?[0-9]*\s?', '', text)
    text = re.sub(r'\\\n', '\n', text)
    text = re.sub(r'\\par\b', '\n', text)
    text = re.sub(r'\\line\b', '\n', text)
    text = re.sub(r'\\tab\b', '\t', text)
    text = re.sub(r'[{}\\]', '', text)
    # Collapse excessive blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


# ─────────────────────────────────────────────────────────────────────────────
#  ODT / ODS / ODP → PDF  (basic extraction)
# ─────────────────────────────────────────────────────────────────────────────

def odf_to_pdf(odf_path: str, pdf_path: str, page_size="a4"):
    """
    Convert ODF formats (ODT, ODS, ODP) to PDF by extracting content.xml
    and rendering text content.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    size = _get_size(page_size)
    story = []
    styles_rl = _make_styles(size)

    NS = {
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
        "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
        "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    }

    try:
        with zipfile.ZipFile(odf_path, "r") as z:
            with z.open("content.xml") as f:
                root = ET.parse(f).getroot()
    except Exception as e:
        story.append(Paragraph(f"Error al leer archivo ODF: {e}", styles_rl["normal"]))
        _build_pdf(story, pdf_path, size)
        return

    body = root.find(".//office:body", NS) or root
    text_el = body.find("office:text", NS) or body

    def get_text(el):
        parts = [el.text or ""]
        for child in el:
            parts.append(get_text(child))
            if child.tail:
                parts.append(child.tail)
        return "".join(parts)

    for el in text_el:
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        text = get_text(el).strip()
        if not text:
            story.append(Spacer(1, 4))
            continue
        text_esc = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        if tag == "h":
            level = int(el.get("{urn:oasis:names:tc:opendocument:xmlns:text:1.0}outline-level", 1))
            s = styles_rl.get(f"h{min(level,4)}", styles_rl["h1"])
            story.append(Paragraph(text_esc, s))
        elif tag == "p":
            story.append(Paragraph(text_esc, styles_rl["normal"]))
        elif tag == "table":
            # Basic table extraction
            rows = []
            for tr in el.findall(".//{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-row"):
                row = []
                for tc in tr.findall(".//{urn:oasis:names:tc:opendocument:xmlns:table:1.0}table-cell"):
                    cell_text = get_text(tc).strip()
                    cell_text = cell_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    row.append(Paragraph(cell_text, styles_rl["table_cell"]))
                if row:
                    rows.append(row)
            if rows:
                t = Table(rows)
                t.setStyle(TableStyle([
                    ("GRID", (0,0),(-1,-1), 0.4, colors.grey),
                    ("TOPPADDING", (0,0),(-1,-1), 3),
                    ("BOTTOMPADDING", (0,0),(-1,-1), 3),
                    ("LEFTPADDING", (0,0),(-1,-1), 5),
                    ("RIGHTPADDING", (0,0),(-1,-1), 5),
                ]))
                story.append(t)
                story.append(Spacer(1, 8))

    _build_pdf(story, pdf_path, size)
    log.info(f"ODF→PDF: {odf_path} → {pdf_path}")


def _build_pdf(story, pdf_path, size):
    doc = SimpleDocTemplate(
        pdf_path, pagesize=size,
        leftMargin=20*mm, rightMargin=20*mm,
        topMargin=20*mm, bottomMargin=20*mm,
    )
    doc.build(story)


# ─────────────────────────────────────────────────────────────────────────────
#  STYLES
# ─────────────────────────────────────────────────────────────────────────────

def _make_styles(page_size):
    base = getSampleStyleSheet()
    W = page_size[0]

    def ps(name, **kw):
        return ParagraphStyle(name, **kw)

    return {
        "normal": ps("pb_normal",
            fontName="Helvetica", fontSize=10, leading=14,
            textColor=colors.HexColor("#1a1a1a"), spaceAfter=4,
            wordWrap="CJK"),
        "title": ps("pb_title",
            fontName="Helvetica-Bold", fontSize=22, leading=28,
            textColor=colors.HexColor("#0f0f2e"), spaceAfter=16,
            alignment=TA_CENTER),
        "h1": ps("pb_h1",
            fontName="Helvetica-Bold", fontSize=17, leading=22,
            textColor=colors.HexColor("#0f3460"), spaceAfter=10,
            spaceBefore=14,
            borderPad=2,),
        "h2": ps("pb_h2",
            fontName="Helvetica-Bold", fontSize=14, leading=18,
            textColor=colors.HexColor("#16213e"), spaceAfter=8,
            spaceBefore=10),
        "h3": ps("pb_h3",
            fontName="Helvetica-Bold", fontSize=12, leading=15,
            textColor=colors.HexColor("#1a1a2e"), spaceAfter=6,
            spaceBefore=8),
        "h4": ps("pb_h4",
            fontName="Helvetica-Bold", fontSize=11, leading=14,
            textColor=colors.HexColor("#333333"), spaceAfter=4),
        "bullet": ps("pb_bullet",
            fontName="Helvetica", fontSize=10, leading=14,
            leftIndent=16, bulletIndent=6,
            textColor=colors.HexColor("#1a1a1a"), spaceAfter=3),
        "mono": ps("pb_mono",
            fontName="Courier", fontSize=8.5, leading=12,
            textColor=colors.HexColor("#222222"), spaceAfter=1),
        "table_cell": ps("pb_cell",
            fontName="Helvetica", fontSize=8, leading=11,
            textColor=colors.HexColor("#111111")),
    }


# ─────────────────────────────────────────────────────────────────────────────
#  MAIN ENTRY POINT
# ─────────────────────────────────────────────────────────────────────────────

CONVERTERS = {
    "docx": docx_to_pdf,
    "doc":  docx_to_pdf,
    "xlsx": xlsx_to_pdf,
    "xls":  xlsx_to_pdf,
    "pptx": pptx_to_pdf,
    "ppt":  pptx_to_pdf,
    "odt":  odf_to_pdf,
    "ods":  odf_to_pdf,
    "odp":  odf_to_pdf,
    "txt":  txt_to_pdf,
    "rtf":  rtf_to_pdf,
}


def _resolve_page_size(page_size: str, orientation: str, src_path: str, ext: str):
    """
    Return (reportlab_page_size, effective_orientation).
    For xlsx/xls: auto-detect if content is wider than portrait → force landscape.
    """
    base = PAGE_SIZES.get(page_size.lower(), A4)

    # XLSX auto-orientation: check column count
    if ext in ("xlsx", "xls") and orientation == "portrait":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(src_path, data_only=True, read_only=True)
            max_col = max((ws.max_column or 1) for ws in wb.worksheets)
            wb.close()
            if max_col > 8:   # more than 8 columns → landscape fits better
                orientation = "landscape"
        except Exception:
            pass

    if orientation == "landscape":
        # reportlab: landscape = (w, h) where w > h
        return (max(base), min(base)), "landscape"
    return base, orientation


def convert_to_pdf(src_path: str, output_dir: str,
                   page_size: str = "a4", orientation: str = "portrait",
                   timeout: int = 120, **kwargs) -> str:
    """
    Convert any supported document to PDF.
    Returns the path to the generated PDF.
    timeout: máximo de segundos para la conversión (default 120).
    """
    import concurrent.futures

    src = Path(src_path)
    ext = src.suffix.lower().lstrip(".")

    converter_fn = CONVERTERS.get(ext)
    if not converter_fn:
        raise RuntimeError(
            f"Formato '.{ext}' no soportado para conversión.\n"
            f"Formatos soportados: {', '.join(f'.{e}' for e in CONVERTERS)}"
        )

    resolved_size, resolved_orient = _resolve_page_size(page_size, orientation, src_path, ext)

    pdf_name = src.stem + ".pdf"
    pdf_path = str(Path(output_dir) / pdf_name)

    def _run():
        extra = {}
        if ext in ("xlsx", "xls"):
            extra["scale"] = kwargs.get("scale", "fit")
        converter_fn(src_path, pdf_path, page_size=resolved_size, **extra)

    # Q5: ejecutar en thread separado con timeout para evitar hangs indefinidos
    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(_run)
        try:
            future.result(timeout=timeout)
        except concurrent.futures.TimeoutError:
            raise RuntimeError(
                f"Timeout ({timeout}s) al convertir {src.name}. "
                "El archivo puede estar corrupto o ser demasiado grande."
            )
        except Exception as e:
            log.exception(f"Conversion failed: {src_path}")
            raise RuntimeError(f"Error al convertir {src.name}: {e}") from e

    if not Path(pdf_path).exists():
        raise RuntimeError(f"La conversión no produjo un PDF: {pdf_path}")

    return pdf_path
